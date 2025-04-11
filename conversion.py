import os
import zipfile
from flask import Flask, request, render_template, send_file
from werkzeug.utils import secure_filename
from pdf2docx import Converter
from docx import Document
import pdfplumber
import fitz  # PyMuPDF
from PIL import Image
import pandas as pd
import PyPDF2
import pytesseract
from fpdf import FPDF
from pptx import Presentation

app = Flask(__name__)

UPLOAD_FOLDER = "uploads"
PROCESSED_FOLDER = "processed"
os.makedirs(UPLOAD_FOLDER, exist_ok=True)
os.makedirs(PROCESSED_FOLDER, exist_ok=True)

app.config["UPLOAD_FOLDER"] = UPLOAD_FOLDER
app.config["PROCESSED_FOLDER"] = PROCESSED_FOLDER
ALLOWED_EXTENSIONS = {"pdf", "docx", "txt", "csv", "html", "jpg", "png", "xlsx", "pptx"}

def allowed_file(filename):
    return "." in filename and filename.rsplit(".", 1)[1].lower() in ALLOWED_EXTENSIONS

def zip_folder(folder_path, zip_name):
    zip_path = os.path.join(PROCESSED_FOLDER, zip_name)
    with zipfile.ZipFile(zip_path, "w") as zipf:
        for root, _, files in os.walk(folder_path):
            for file in files:
                file_path = os.path.join(root, file)
                arcname = os.path.relpath(file_path, folder_path)
                zipf.write(file_path, arcname)
    return zip_path

@app.route("/", methods=["GET", "POST"])
def upload_file():
    if request.method == "POST":
        files = request.files.getlist("file")
        action = request.form.get("action")

        if not files:
            return "No file selected"

        file_paths = []
        for file in files:
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                file_path = os.path.join(UPLOAD_FOLDER, filename)
                file.save(file_path)
                file_paths.append(file_path)

        file_path = file_paths[0]
        output_file = ""

        if action == "pdf_to_word":
            output_file = file_path.replace(".pdf", ".docx")
            cv = Converter(file_path)
            cv.convert(output_file)
            cv.close()

        elif action == "word_to_pdf":
            output_file = file_path.replace(".docx", ".pdf")
            doc = Document(file_path)
            doc.save(output_file)

        elif action == "txt_to_word":
            output_file = file_path.replace(".txt", ".docx")
            doc = Document()
            with open(file_path, "r", encoding="utf-8") as f:
                doc.add_paragraph(f.read())
            doc.save(output_file)

        elif action == "txt_to_pdf":
            output_file = file_path.replace(".txt", ".pdf")
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            with open(file_path, "r", encoding="utf-8") as f:
                for line in f:
                    pdf.cell(200, 10, line, ln=True)
            pdf.output(output_file)

        elif action == "excel_to_pdf":
            df = pd.read_excel(file_path)
            output_file = file_path.replace(".xlsx", ".pdf")
            df.to_csv("temp.csv", index=False)
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            with open("temp.csv", "r", encoding="utf-8") as f:
                for line in f:
                    pdf.cell(200, 10, line, ln=True)
            pdf.output(output_file)

        elif action == "pptx_to_pdf":
            prs = Presentation(file_path)
            output_file = file_path.replace(".pptx", ".pdf")
            pdf = FPDF()
            pdf.set_font("Arial", size=12)
            pdf.add_page()
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        pdf.cell(200, 10, shape.text, ln=True)
            pdf.output(output_file)

        elif action == "image_to_text":
            output_file = file_path.replace(".png", ".txt").replace(".jpg", ".txt")
            text = pytesseract.image_to_string(Image.open(file_path))
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(text)

        elif action == "csv_to_excel":
            output_file = file_path.replace(".csv", ".xlsx")
            df = pd.read_csv(file_path)
            df.to_excel(output_file, index=False)

        elif action == "excel_to_csv":
            output_file = file_path.replace(".xlsx", ".csv")
            df = pd.read_excel(file_path)
            df.to_csv(output_file, index=False)

        elif action == "html_to_pdf":
            output_file = file_path.replace(".html", ".pdf")
            with open(file_path, "r", encoding="utf-8") as f:
                html_content = f.read()
            pdf = FPDF()
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.multi_cell(0, 10, html_content)
            pdf.output(output_file)

        elif action == "pdf_merge":
            merger = PyPDF2.PdfMerger()
            for path in file_paths:
                merger.append(path)
            output_file = os.path.join(PROCESSED_FOLDER, "merged.pdf")
            merger.write(output_file)
            merger.close()

        elif action == "pdf_to_txt":
            with pdfplumber.open(file_path) as pdf:
                text = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
            output_file = file_path.replace(".pdf", ".txt")
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(text)

        elif action == "pdf_to_html":
            with pdfplumber.open(file_path) as pdf:
                text = "\n".join([page.extract_text() for page in pdf.pages if page.extract_text()])
            output_file = file_path.replace(".pdf", ".html")
            with open(output_file, "w", encoding="utf-8") as f:
                f.write(f"<html><body><pre>{text}</pre></body></html>")

        elif action == "pdf_to_excel":
            output_file = file_path.replace(".pdf", ".xlsx")
            df = pd.read_csv(file_path)  # Placeholder logic
            df.to_excel(output_file, index=False)

        elif action == "pdf_to_pptx":
            output_file = file_path.replace(".pdf", ".pptx")
            prs = Presentation()
            slide_layout = prs.slide_layouts[5]
            with pdfplumber.open(file_path) as pdf:
                for page in pdf.pages:
                    slide = prs.slides.add_slide(slide_layout)
                    textbox = slide.shapes.add_textbox(10, 10, 600, 500)
                    textbox.text = page.extract_text()
            prs.save(output_file)

        elif action == "image_to_pdf":
            images = [Image.open(f).convert("RGB") for f in file_paths]
            output_file = os.path.join(PROCESSED_FOLDER, "converted_images.pdf")
            images[0].save(output_file, save_all=True, append_images=images[1:])

        elif action == "pdf_to_image":
            doc = fitz.open(file_path)
            output_folder = os.path.join(PROCESSED_FOLDER, "pdf_images")
            os.makedirs(output_folder, exist_ok=True)
            image_paths = []
            for i, page in enumerate(doc):
                pix = page.get_pixmap()
                img_path = os.path.join(output_folder, f"page_{i + 1}.png")
                pix.save(img_path)
                image_paths.append(img_path)
            output_file = zip_folder(output_folder, "pdf_images.zip")

        elif action == "add_watermark":
            output_file = os.path.join(PROCESSED_FOLDER, "watermarked.pdf")
            watermark_text = "CONFIDENTIAL"
            pdf = FPDF()
            pdf.set_auto_page_break(auto=True, margin=15)
            pdf.add_page()
            pdf.set_font("Arial", size=12)
            pdf.set_xy(50, 50)
            pdf.cell(0, 10, watermark_text, 0, 1, "C")
            pdf.output(output_file)

        elif action == "pdf_split":
            output_folder = os.path.join(PROCESSED_FOLDER, "split_pages")
            os.makedirs(output_folder, exist_ok=True)
            pdf_reader = PyPDF2.PdfReader(file_path)
            for i, page in enumerate(pdf_reader.pages):
                pdf_writer = PyPDF2.PdfWriter()
                pdf_writer.add_page(page)
                split_output = os.path.join(output_folder, f"split_page_{i+1}.pdf")
                with open(split_output, "wb") as output_pdf:
                    pdf_writer.write(output_pdf)
            output_file = zip_folder(output_folder, "pdf_split_pages.zip")

        if output_file and os.path.exists(output_file):
            return send_file(output_file, as_attachment=True)
        else:
            return "Conversion failed or not implemented."

    return render_template("index.html")

if __name__ == "__main__":
    app.run(debug=True)


